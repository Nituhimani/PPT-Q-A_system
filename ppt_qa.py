import os
from pptx import Presentation
from PIL import Image
import pytesseract
import cv2
import numpy as np
from transformers import pipeline
import tempfile
import io

class PPTProcessor:
    def __init__(self):
        # Initialize the question-answering pipeline
        self.qa_pipeline = pipeline("question-answering")
        self.content = []
        
    def extract_text_from_slide(self, slide):
        """Extract text from a slide."""
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        return " ".join(text_content)
    
    def extract_image_from_slide(self, slide):
        """Extract and process images from a slide."""
        image_content = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                image_bytes = image.blob
                
                # Convert to PIL Image
                image = Image.open(io.BytesIO(image_bytes))
                
                # Convert to OpenCV format for processing
                cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
                
                # Extract text from image using OCR
                text = pytesseract.image_to_string(cv_image)
                if text.strip():
                    image_content.append(f"Image content: {text}")
                
                # TODO: Add more sophisticated image analysis here
                # This could include object detection, graph analysis, etc.
                
        return image_content
    
    def process_ppt(self, ppt_path):
        """Process the entire PowerPoint file."""
        prs = Presentation(ppt_path)
        
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_content = {
                'slide_number': slide_number,
                'text': self.extract_text_from_slide(slide),
                'images': self.extract_image_from_slide(slide)
            }
            self.content.append(slide_content)
    
    def answer_question(self, question):
        """Answer a question based on the processed content."""
        # Combine all content into a single context
        context = ""
        for slide in self.content:
            context += f"Slide {slide['slide_number']}: {slide['text']}\n"
            for image_content in slide['images']:
                context += f"{image_content}\n"
        
        # Use the QA pipeline to get the answer
        result = self.qa_pipeline(
            question=question,
            context=context
        )
        
        return result['answer']

def main():
    processor = PPTProcessor()
    
    # Get the PPT file path from user
    ppt_path = input("Enter your ppt file path: ")
    
    if not os.path.exists(ppt_path):
        print("File not found!")
        return
    
    print("Processing PowerPoint file...")
    processor.process_ppt(ppt_path)
    print("Processing complete!")
    
    while True:
        question = input("\nEnter your question (or 'quit' to exit): ")
        if question.lower() == 'quit':
            break
            
        try:
            answer = processor.answer_question(question)
            print(f"\nAnswer: {answer}")
        except Exception as e:
            print(f"Error processing question: {str(e)}")

if __name__ == "__main__":
    main() 